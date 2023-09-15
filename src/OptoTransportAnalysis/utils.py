from tkinter import filedialog as fd
import os
import json
import pptx
from pptx.util import Inches
from cycler import cycler


#### Utility variables -------------------------------------------------------

## Dimensions of relevant quantities in a powerpoint
prs_width = Inches(13.333)
prs_height = Inches(7.5)
pt18_height = Inches(1/4)

## Default plotting parameters to use
default_rcParams = {'axes.labelsize' : 14,
                    'axes.titlesize' : 16,
                    'axes.titleweight' : 2,
                    'axes.formatter.use_mathtext' : False,
                    'axes.prop_cycle' : cycler(color=['red', 'orange', 'gold', 'green', 'blue', 'indigo', 'violet']),
                    'axes.labelpad' : 10,
                    'lines.linewidth' : 0.8,
                    'legend.fontsize' : 8,
                    'legend.fancybox' : False,
                    'legend.shadow' : False,
                    'legend.title_fontsize' : 14,
                    'savefig.dpi' : 1200,
                    'savefig.transparent' : True,
                    'savefig.format' : 'png',
                    'savefig.bbox' : 'tight',
                    'figure.dpi' : 225,
                    'figure.figsize' : (6,4),
                    'xtick.labelsize' : 12,
                    'ytick.labelsize' : 12}


#### Utility functions -------------------------------------------------------


def create_metadata_from_dict(dict_to_write, initial_dir=""):
    """
    Create a .json metadata file from a dict.

    Useful if given data from someone else, if forgot to record a 
    metadata file when first acquiring data, or if associated 
    metadata file is lost.

    Parameters
    ----------
    dict_to_write : dict
        Dict being written to .json file

    initial_dir : str, default=""
        Directory to begin search from. On function call, a 
        tkinter.filedialog GUI is opened to select existing data file 
        to associate the new metadata file with. 

    Returns
    -------

    N/A

    """
    # Prompt user to select file with tkinter.filedialog GUI
    prompt_str = "Please select the file you wish to generate metadata for"
    filename = fd.askopenfilename(title=prompt_str, initialdir=initial_dir)

    # Solve for .json filename and dump to .json format
    filename_json = os.path.splitext(filename)[0] + '.json'
    with open(filename_json, "w") as file:
        json.dump(dict_to_write, file)


def create_summary_slide(prs, fig_name, run_title=None, 
                         custom_title=None, comment_txt=None):
    """
    Add a brief summary slide containing a figure and comment to an 
    existing powerpoint presentation.

    Parameters
    ----------

    prs : pptx.Presentation
        Existing PowerPoint presentation to edit.

    fig_name : path, path-like str
        Directory at which the figure to be put in the slide can be found.

    run_title : str, optional
        The title of the run in the Data object from which data is 
        being plotted. If given, sets the title of the slide to contain 
        the name of the data file.

    custom_title : str, optional
        If run_title not given, can also give a custom title. Useful if 
        plotted data comes from multiple datasets, etc. If both run_title
        and custom_title are given, custom_title takes precedence.

    comment_txt : str, optional
        Optional comment which can be added to the slide.

    Returns
    -------

    slide : ?
        Unsure of the exact type of this object, but seems to be a 
        reference to the created slide such that edits can be made 
        after creation.

    title : ?
        Also unsure of the exact type of this object, but seems to be a
        reference to the title of the slide such that edits can be made
        after creation.

    figure : ?
        Also unsure of the exact type of this object, but seems to be a
        reference to the figure of the slide such that edits can be made
        after creation.
    
    comment : ?
        Also unsure of the exact type of this object, but seems to be a
        reference to the comment of the slide such that edits can be made
        after creation.

    Notes
    -----

    For more information on the python-pptx package, look to
    https://python-pptx.readthedocs.io/en/latest/index.html

    """
    # Create slide using blank layout option
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Give the slide a title
    # -- If run_title given and custom_title not given, set run_title as title
    # -- If custom_title given, set custom_title as title
    # -- Else, add default title to be edited later
    title = slide.shapes.add_textbox(left=0, top=0, width=prs_width, height=2*pt18_height)
    if run_title != None and custom_title == None:
        title.text = run_title
    elif (run_title == None and custom_title != None) or (run_title != None and custom_title != None):
        title.text = custom_title
    else:
        title.text = "Add title"

    # Add a figure to the slide and center it
    figure = slide.shapes.add_picture(image_file=fig_name, top=title.height, 
                                      left=0, height=prs_height - 3*title.height)
    figure.left = (prs_width - figure.width) // 2

    # Add a comment which can be edited later
    comment = slide.shapes.add_textbox(left=0, top=figure.height + title.height, width=prs_width, height=2*pt18_height)
    if comment != None:
        comment.text = comment_txt

    return [slide, title, figure, comment]

###### Plotting utility functions --------------------------------------------

def label_H_x(ax):
    """
    Names the x_label of a plot after the applied magnetic field in Tesla.

    Parameters
    ----------

    ax : Axes
        Axes object you wish to label the x-axis of.

    Returns
    -------

    N/A

    """
    ax.set_xlabel(r'μ_0H (T)')
    return


def label_IDC_uA_x(ax):
    """
    Names the x_label of a plot after the DC offset current in μA.

    Parameters
    ----------

    ax : Axes
        Axes object you wish to label the x-axis of.

    Returns
    -------
    N/A
    
    """
    ax.set_xlabel(r'I$_{DC}$ (μA)')
    return


def label_photon_wavelength_nm_X(ax):
    """
    Names the x_label of a plot after the collection wavelength in nm.

    Parameters
    ----------

    ax : Axes
        Axes object you wish to label the x-axis of.

    Returns
    -------
    N/A
    
    """
    ax.set_xlabel(r'λ (nm)')
    return


def label_photon_energy_eV_x(ax):
    """
    Names the x_label of a plot after the collection energy in eV.

    Parameters
    ----------

    ax : Axes
        Axes object you wish to label the x-axis of.

    Returns
    -------
    N/A
    
    """
    ax.set_xlabel(r'$\mathcal{E}$ (eV)')
    return